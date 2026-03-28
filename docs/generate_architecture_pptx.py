"""
Generate ForexMind Architecture PowerPoint
==========================================
Run:  python docs/generate_architecture_pptx.py
Output: docs/ForexMind_Architecture.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm

OUTPUT = os.path.join(os.path.dirname(__file__), "ForexMind_Architecture.pptx")

# ── Colour palette ────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0D, 0x1B, 0x2A)
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
TEAL   = RGBColor(0x00, 0xBF, 0xA5)
GOLD   = RGBColor(0xF9, 0xA8, 0x25)
LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
GREY   = RGBColor(0xF5, 0xF5, 0xF5)
MID    = RGBColor(0x90, 0xA4, 0xAE)
GREEN  = RGBColor(0x43, 0xA0, 0x47)
RED    = RGBColor(0xE5, 0x39, 0x35)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x21, 0x21, 0x21)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Slide helpers ─────────────────────────────────────────────────────────────

def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, l, t, w, h)  # 1 = MSO_SHAPE_TYPE.RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, l, t, w, h, font_size=18, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, italic=False,
                 font_name="Calibri"):
    txb = slide.shapes.add_textbox(l, t, w, h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txb


def add_table(slide, headers, rows, l, t, w, h,
              header_bg=BLUE, header_fg=WHITE,
              row_bgs=None, font_size=11):
    if row_bgs is None:
        row_bgs = [WHITE, GREY]
    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl = slide.shapes.add_table(num_rows, num_cols, l, t, w, h).table

    # Column widths — equal split
    col_w = w // num_cols
    for i in range(num_cols):
        tbl.columns[i].width = col_w

    def _cell(cell, text, bg, fg, bold=False):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = fg
        run.font.name = "Calibri"

    for ci, hdr in enumerate(headers):
        _cell(tbl.cell(0, ci), hdr, header_bg, header_fg, bold=True)

    for ri, row in enumerate(rows):
        bg = row_bgs[ri % len(row_bgs)]
        for ci, val in enumerate(row):
            _cell(tbl.cell(ri + 1, ci), str(val), bg, DARK)

    return tbl


def add_bullet_box(slide, title, bullets, l, t, w, h,
                   title_color=GOLD, bullet_color=WHITE, bg_color=NAVY):
    # Background
    rect = add_rect(slide, l, t, w, h, bg_color)

    # Title
    add_text_box(slide, title, l + Inches(0.15), t + Inches(0.1),
                 w - Inches(0.3), Inches(0.45),
                 font_size=14, bold=True, color=title_color)

    # Bullets
    txb = slide.shapes.add_textbox(
        l + Inches(0.25), t + Inches(0.6), w - Inches(0.4), h - Inches(0.7)
    )
    tf = txb.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"• {b}"
        run.font.size = Pt(11)
        run.font.color.rgb = bullet_color
        run.font.name = "Calibri"
        p.space_after = Pt(4)


def title_bar(slide, title, subtitle=""):
    # Dark header bar
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), NAVY)
    add_text_box(slide, title, Inches(0.3), Inches(0.1), Inches(12.5), Inches(0.6),
                 font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle, Inches(0.3), Inches(0.68), Inches(12.5), Inches(0.35),
                     font_size=13, color=GOLD)


# ── Slide 1: Cover ────────────────────────────────────────────────────────────

def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, NAVY)

    # Gold accent bar top
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), GOLD)

    # Title
    add_text_box(slide, "ForexMind",
                 Inches(1), Inches(1.4), Inches(11.33), Inches(1.5),
                 font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide, "AI-Powered Forex Trading Agent",
                 Inches(1), Inches(2.8), Inches(11.33), Inches(0.7),
                 font_size=24, color=LIGHT, align=PP_ALIGN.CENTER)

    # Tagline
    add_text_box(slide,
                 "System Architecture & Technical Design  |  Version 1.0  |  March 2026",
                 Inches(1), Inches(3.4), Inches(11.33), Inches(0.5),
                 font_size=14, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

    # Stats boxes
    stats = [
        ("≥70%", "Target\nAccuracy"),
        ("4",    "ML\nStrategies"),
        ("18",   "Forex\nPairs"),
        ("3",    "User\nInterfaces"),
    ]
    box_w = Inches(2.8)
    for i, (val, label) in enumerate(stats):
        lx = Inches(0.5 + i * 3.1)
        add_rect(slide, lx, Inches(4.3), box_w, Inches(1.7), BLUE)
        add_text_box(slide, val, lx, Inches(4.4), box_w, Inches(0.8),
                     font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, lx, Inches(5.1), box_w, Inches(0.7),
                     font_size=12, color=LIGHT, align=PP_ALIGN.CENTER)

    # Gold accent bar bottom
    add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), GOLD)


# ── Slide 2: Agenda ───────────────────────────────────────────────────────────

def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, GREY)
    title_bar(slide, "Agenda", "What this presentation covers")

    items = [
        ("1", "System Overview & Design Philosophy"),
        ("2", "High-Level Architecture"),
        ("3", "Data Layer — OANDA & News"),
        ("4", "Technical Indicators Engine"),
        ("5", "Four-Strategy Ensemble"),
        ("6", "Risk Management"),
        ("7", "Claude AI Agent (LangChain 1.x)"),
        ("8", "Backtesting & Validation"),
        ("9", "User Interfaces"),
        ("10", "Configuration & Security"),
        ("11", "Performance Targets"),
    ]

    col1 = items[:6]
    col2 = items[6:]

    for col_idx, col_items in enumerate([col1, col2]):
        lx = Inches(0.4 + col_idx * 6.5)
        for row_idx, (num, text) in enumerate(col_items):
            ty = Inches(1.3 + row_idx * 0.87)
            add_rect(slide, lx, ty, Inches(0.5), Inches(0.6), BLUE)
            add_text_box(slide, num, lx, ty + Inches(0.05), Inches(0.5), Inches(0.5),
                         font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_rect(slide, lx + Inches(0.55), ty, Inches(5.7), Inches(0.6), WHITE)
            add_text_box(slide, text, lx + Inches(0.65), ty + Inches(0.08),
                         Inches(5.5), Inches(0.5), font_size=13, color=DARK)


# ── Slide 3: Overview ─────────────────────────────────────────────────────────

def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "System Overview", "What ForexMind does and why it works")

    add_bullet_box(slide, "Core Mission",
        ["Achieve ≥70% directional accuracy on 1–5 minute forex charts",
         "All major, minor and exotic pairs (18 total)",
         "Conversational AI interface — ask questions in plain English",
         "AI-decided position sizing with hard risk controls"],
        Inches(0.3), Inches(1.2), Inches(6.1), Inches(2.8))

    add_bullet_box(slide, "Design Principles",
        ["Ensemble over single-model: 4 orthogonal strategies vote",
         "Graceful degradation: untrained models return HOLD",
         "Risk-first: no signal bypasses RiskManager",
         "Async-native: all I/O non-blocking",
         "Backtest before live: statistical edge required"],
        Inches(6.7), Inches(1.2), Inches(6.3), Inches(2.8))

    # Technology pills
    tech = ["Claude 3.5 Sonnet", "LangChain 1.x", "LightGBM", "PyTorch BiLSTM",
            "PPO-RL", "pandas-ta", "OANDA REST", "FastAPI", "Rich CLI", "Telegram Bot"]
    for i, t in enumerate(tech):
        lx = Inches(0.3 + (i % 5) * 2.55)
        ty = Inches(4.3 + (i // 5) * 0.6)
        add_rect(slide, lx, ty, Inches(2.3), Inches(0.45), TEAL)
        add_text_box(slide, t, lx, ty + Inches(0.04), Inches(2.3), Inches(0.4),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ── Slide 4: Architecture Diagram ────────────────────────────────────────────

def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK)
    title_bar(slide, "High-Level Architecture", "End-to-end data flow")

    layers = [
        (GREY,  DARK,  "EXTERNAL SOURCES",     "OANDA REST API  •  Alpha Vantage News  •  Finnhub News"),
        (BLUE,  WHITE, "DATA LAYER",            "OandaClient (async)  •  NewsAggregator (dedup + sentiment)  •  SQLite"),
        (TEAL,  DARK,  "INDICATORS (25+)",      "EMA / MACD / RSI / ADX / Stochastic / Bollinger / ATR / PSAR / Pivots ..."),
        (GREEN, WHITE, "SIGNAL SCORER",         "Weighted composite score  [ -100 → +100 ]   Trend 30% · Momentum 25% · Structure 20%"),
        (BLUE,  WHITE, "ENSEMBLE STRATEGY",     "Rule-Based 30%  |  LightGBM 25%  |  Bi-LSTM 25%  |  PPO-RL 20%"),
        (GOLD,  DARK,  "RISK MANAGER",          "Kill-switch · ATR stops · Kelly sizing · R:R ≥ 1:1 · Trailing stops"),
        (RGBColor(0x6A,0x1B,0x9A), WHITE, "CLAUDE AI AGENT", "LangChain 1.x · 6 tools · Sliding window memory · JSON signal synthesis"),
        (GREY,  DARK,  "USER INTERFACES",       "Rich CLI  •  FastAPI Dashboard + WebSocket  •  Telegram Bot"),
    ]

    box_h = Inches(0.68)
    for i, (bg, fg, label, detail) in enumerate(layers):
        ty = Inches(1.15 + i * 0.76)
        add_rect(slide, Inches(0.3), ty, Inches(2.5), box_h, bg)
        add_text_box(slide, label, Inches(0.3), ty + Inches(0.12),
                     Inches(2.5), box_h - Inches(0.24),
                     font_size=11, bold=True, color=fg, align=PP_ALIGN.CENTER)
        add_rect(slide, Inches(2.85), ty, Inches(9.9), box_h, RGBColor(0x2A, 0x2A, 0x3A))
        add_text_box(slide, detail, Inches(3.0), ty + Inches(0.14),
                     Inches(9.6), box_h - Inches(0.28),
                     font_size=11, color=LIGHT)

        if i < len(layers) - 1:
            add_text_box(slide, "▼", Inches(0.3), ty + box_h, Inches(2.5), Inches(0.18),
                         font_size=10, color=MID, align=PP_ALIGN.CENTER)


# ── Slide 5: Data Layer ───────────────────────────────────────────────────────

def slide_data_layer(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Data Layer", "OANDA Client & News Aggregator")

    add_bullet_box(slide, "OANDA REST Client",
        ["asyncio.to_thread() wraps sync oandapyV20 library — fully non-blocking",
         "@retry decorator: 3 attempts, exponential back-off (×2)",
         "get_candles() → pandas DataFrame with UTC DatetimeIndex",
         "get_multi_candles() → concurrent fetch for multiple pairs",
         "market_order() with SL/TP bracket on fill",
         "Targets: api-fxpractice.oanda.com (free, $100k virtual)"],
        Inches(0.3), Inches(1.2), Inches(6.2), Inches(3.2), title_color=GOLD)

    add_bullet_box(slide, "News Aggregator",
        ["NewsSource ABC — 2 implementations: AlphaVantage + Finnhub",
         "Deduplication by headline hash (no repeated articles)",
         "15-minute TTL cache — respects free API rate limits",
         "TextBlob sentiment: -1.0 (bearish) → +1.0 (bullish)",
         "get_instrument_sentiment() → score + impact + headlines",
         "Free tiers: AV 25/day, Finnhub 60/min"],
        Inches(6.8), Inches(1.2), Inches(6.2), Inches(3.2), title_color=GOLD)

    # Session detector box
    add_rect(slide, Inches(0.3), Inches(4.6), Inches(12.7), Inches(1.5), LIGHT)
    add_text_box(slide, "Session Times Detector (utils/session_times.py)",
                 Inches(0.5), Inches(4.65), Inches(12.3), Inches(0.4),
                 font_size=13, bold=True, color=NAVY)
    sessions = "Sydney 21:00–06:00  •  Tokyo 00:00–09:00  •  London 07:00–16:00  •  New York 12:00–21:00  •  Prime overlaps: London/NY 12:00–16:00"
    add_text_box(slide, sessions, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.5),
                 font_size=11, color=DARK)


# ── Slide 6: Indicators ───────────────────────────────────────────────────────

def slide_indicators(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Technical Indicators Engine", "25+ indicators via pandas-ta (pure Python)")

    categories = [
        ("Trend",      NAVY,  "EMA 9/21/50/200  •  PSAR",            "ema_trend, psar_signal"),
        ("Momentum",   BLUE,  "MACD(12/26/9)  •  RSI(14)  •  Stoch(14,3)", "macd_cross, rsi_zone, stoch_cross"),
        ("Oscillators",TEAL,  "CCI(20)  •  Williams %R(14)  •  MFI(14)",  "Position vs extremes"),
        ("Trend Str.", GREEN, "ADX(14) + DI+ / DI−",                 "adx_trend_strength"),
        ("Volatility", GOLD,  "ATR(14)  •  Bollinger Bands(20,2)",   "bb_position, atr for stop sizing"),
        ("Volume",     MID,   "OBV  •  ROC(5/10)",                   "Volume trend confirmation"),
        ("Structure",  RED,   "Pivot Points  •  Swing highs/lows",   "support, resistance levels"),
    ]

    for i, (cat, col, inds, signal) in enumerate(categories):
        row = i // 2
        col_pos = i % 2
        lx = Inches(0.3 + col_pos * 6.55)
        ty = Inches(1.25 + row * 1.55)
        bw = Inches(6.1)
        bh = Inches(1.35)
        add_rect(slide, lx, ty, bw, bh, col)
        add_text_box(slide, cat, lx + Inches(0.1), ty + Inches(0.07),
                     bw - Inches(0.2), Inches(0.35),
                     font_size=13, bold=True, color=WHITE)
        add_text_box(slide, inds, lx + Inches(0.1), ty + Inches(0.4),
                     bw - Inches(0.2), Inches(0.4),
                     font_size=10, color=WHITE)
        add_text_box(slide, f"→ {signal}", lx + Inches(0.1), ty + Inches(0.82),
                     bw - Inches(0.2), Inches(0.35),
                     font_size=10, color=GOLD, italic=True)

    # Last box (7th) is solo
    # Signal Scorer summary
    add_rect(slide, Inches(0.3), Inches(6.15), Inches(12.7), Inches(0.9), LIGHT)
    add_text_box(slide,
                 "Signal Scorer: Trend 30%  •  Momentum 25%  •  Structure 20%  •  Volume 15%  •  Volatility 10%  →  Score [-100, +100]",
                 Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.7),
                 font_size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ── Slide 7: Ensemble ─────────────────────────────────────────────────────────

def slide_ensemble(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Four-Strategy Ensemble", "Weighted voting for robust directional prediction")

    strategies = [
        (NAVY,  "Rule-Based\n30%",   "6-condition checklist\n(HTF trend, EMA stack,\nMACD, RSI, Stoch, PSAR)\nConfidence = conditions/6"),
        (BLUE,  "LightGBM\n25%",     "Gradient-boosted tree\nTimeSeriesSplit 5-fold CV\nFeature lags + rolling stats\nbalanced class weights"),
        (TEAL,  "BiLSTM\n25%",       "2-layer BiLSTM (128 hidden)\nAdamW + CosineAnnealing\nBatchNorm + dropout=0.3\n3-class softmax output"),
        (GREEN, "PPO-RL\n20%",       "stable-baselines3 PPO\nCustom gymnasium.Env\nDiscrete(3) actions\nP&L reward − overtrading"),
    ]

    for i, (col, title, desc) in enumerate(strategies):
        lx = Inches(0.3 + i * 3.25)
        add_rect(slide, lx, Inches(1.2), Inches(3.0), Inches(2.1), col)
        add_text_box(slide, title, lx, Inches(1.3),
                     Inches(3.0), Inches(0.75),
                     font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, lx + Inches(0.1), Inches(2.05),
                     Inches(2.8), Inches(1.15),
                     font_size=10, color=WHITE)

    # Voting algorithm
    add_rect(slide, Inches(0.3), Inches(3.5), Inches(12.7), Inches(1.3), LIGHT)
    add_text_box(slide, "Ensemble Voting Algorithm",
                 Inches(0.5), Inches(3.55), Inches(12.0), Inches(0.38),
                 font_size=13, bold=True, color=NAVY)
    add_text_box(slide,
                 "buy_score = Σ(weight × confidence)  for BUY strategies      "
                 "sell_score = Σ(weight × confidence)  for SELL strategies\n"
                 "Direction wins if score ≥ 0.50  AND  ≥2 strategies agree   →   otherwise HOLD",
                 Inches(0.5), Inches(3.93), Inches(12.0), Inches(0.8),
                 font_size=11, color=DARK)

    # Feature engineering summary
    add_rect(slide, Inches(0.3), Inches(5.0), Inches(12.7), Inches(1.7), NAVY)
    add_text_box(slide, "Feature Engineering Pipeline (strategy/feature_engineering.py)",
                 Inches(0.5), Inches(5.07), Inches(12.0), Inches(0.4),
                 font_size=12, bold=True, color=GOLD)
    feats = ("5-bar lag features  •  Rolling stats (windows 5/10/20)  •  Candle patterns (engulfing, doji, body%)  "
             "•  Session flags  •  Normalised indicators  •  Target label: next-bar direction ≥3 pip threshold")
    add_text_box(slide, feats, Inches(0.5), Inches(5.5), Inches(12.0), Inches(1.0),
                 font_size=11, color=LIGHT)


# ── Slide 8: Risk Management ──────────────────────────────────────────────────

def slide_risk(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Risk Management Engine", "The gatekeeper — no trade can bypass it")

    # Pipeline flow
    steps = [
        (RED,   "1. Kill-Switch",    "Daily loss > −5%\n→ ALL trading halted"),
        (GOLD,  "2. Concurrent",     "≥5 open trades\n→ Reject"),
        (BLUE,  "3. ATR Stops",      "SL = entry ± ATR×1.5\nTP = SL×RR_ratio"),
        (TEAL,  "4. R:R Check",      "tp_pips / sl_pips\n≥ 1.0  or  Reject"),
        (GREEN, "5. Kelly Size",     "risk_pct = kelly×50\ncapped at 0.5%–2.0%"),
        (NAVY,  "6. RiskProposal",   "approved=True/False\nwith rejection reason"),
    ]

    for i, (col, title, desc) in enumerate(steps):
        lx = Inches(0.3 + i * 2.17)
        add_rect(slide, lx, Inches(1.2), Inches(2.0), Inches(2.2), col)
        add_text_box(slide, title, lx, Inches(1.28),
                     Inches(2.0), Inches(0.55),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, lx + Inches(0.05), Inches(1.85),
                     Inches(1.9), Inches(1.4),
                     font_size=10, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(steps) - 1:
            add_text_box(slide, "→", Inches(2.27 + i * 2.17), Inches(1.9),
                         Inches(0.22), Inches(0.6),
                         font_size=18, bold=True, color=MID, align=PP_ALIGN.CENTER)

    # Kelly formula
    add_rect(slide, Inches(0.3), Inches(3.6), Inches(6.1), Inches(1.5), LIGHT)
    add_text_box(slide, "Kelly Criterion (Half-Kelly)",
                 Inches(0.5), Inches(3.65), Inches(5.7), Inches(0.4),
                 font_size=12, bold=True, color=NAVY)
    add_text_box(slide,
                 "kelly = (win_rate × rr_ratio − (1 − win_rate)) / rr_ratio\n"
                 "risk_pct = clamp(kelly × 50,  min=0.5%,  max=2.0%)",
                 Inches(0.5), Inches(4.05), Inches(5.7), Inches(0.9),
                 font_size=11, color=DARK)

    # Trailing stop
    add_rect(slide, Inches(6.7), Inches(3.6), Inches(6.3), Inches(1.5), LIGHT)
    add_text_box(slide, "Live Trade Management",
                 Inches(6.9), Inches(3.65), Inches(5.9), Inches(0.4),
                 font_size=12, bold=True, color=NAVY)
    add_text_box(slide,
                 "• Trailing stop: moves with price at ATR×1.0 distance\n"
                 "• Break-even escalation after 1R profit\n"
                 "• asyncio.Lock ensures thread-safe state updates",
                 Inches(6.9), Inches(4.05), Inches(5.9), Inches(0.9),
                 font_size=11, color=DARK)

    # Spread filter note
    add_rect(slide, Inches(0.3), Inches(5.3), Inches(12.7), Inches(0.8), NAVY)
    add_text_box(slide,
                 "Spread Filter: skip trade if bid–ask spread > 3 pips  •  "
                 "Paper Trading mode ON by default (PAPER_TRADING=true in .env)  •  "
                 "Max concurrent trades: 5",
                 Inches(0.5), Inches(5.38), Inches(12.3), Inches(0.6),
                 font_size=11, color=LIGHT, align=PP_ALIGN.CENTER)


# ── Slide 9: Claude Agent ─────────────────────────────────────────────────────

def slide_agent(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Claude AI Agent", "LangChain 1.x  bind_tools  +  manual agentic loop")

    # Agentic loop diagram
    loop_steps = [
        (NAVY, "User\nMessage"),
        (BLUE, "Build\nMessages"),
        (TEAL, "LLM\nInvoke"),
        (GOLD, "Tool\nCalls?"),
        (GREEN,"Execute\nTools"),
        (NAVY, "Final\nAnswer"),
    ]
    for i, (col, label) in enumerate(loop_steps):
        lx = Inches(0.3 + i * 2.1)
        add_rect(slide, lx, Inches(1.2), Inches(1.85), Inches(0.9), col)
        add_text_box(slide, label, lx, Inches(1.25),
                     Inches(1.85), Inches(0.8),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        if i < len(loop_steps) - 1:
            arrow = "→" if i != 3 else "↓/→"
            add_text_box(slide, arrow, Inches(2.1 + i * 2.1), Inches(1.4),
                         Inches(0.28), Inches(0.5),
                         font_size=16, bold=True, color=MID, align=PP_ALIGN.CENTER)

    add_text_box(slide, "loop back if tool calls exist (max 6 iterations)",
                 Inches(3.5), Inches(2.15), Inches(6), Inches(0.35),
                 font_size=9, color=MID, italic=True)

    # Six tools
    tools = [
        (NAVY,  "get_signal",    "Live OANDA data + indicators + ensemble → JSON signal"),
        (BLUE,  "get_news",      "Alpha Vantage + Finnhub → sentiment summary"),
        (TEAL,  "get_account",   "Balance, NAV, margin, daily P&L, kill-switch status"),
        (GREEN, "run_backtest",  "Historical data → Sharpe, win rate, profit factor"),
        (GOLD,  "place_trade",   "Market order with SL/TP on OANDA paper account"),
        (MID,   "get_sessions",  "Active sessions, overlaps, recommended pairs"),
    ]

    add_text_box(slide, "6 LangChain Tools (StructuredTool with Pydantic schemas)",
                 Inches(0.3), Inches(2.65), Inches(12.7), Inches(0.38),
                 font_size=12, bold=True, color=NAVY)

    for i, (col, name, desc) in enumerate(tools):
        row = i // 3
        col_pos = i % 3
        lx = Inches(0.3 + col_pos * 4.35)
        ty = Inches(3.1 + row * 0.85)
        add_rect(slide, lx, ty, Inches(1.4), Inches(0.7), col)
        add_text_box(slide, name, lx, ty + Inches(0.05),
                     Inches(1.4), Inches(0.6),
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, lx + Inches(1.45), ty, Inches(2.8), Inches(0.7), LIGHT)
        add_text_box(slide, desc, lx + Inches(1.55), ty + Inches(0.1),
                     Inches(2.6), Inches(0.55),
                     font_size=10, color=DARK)

    # Memory
    add_rect(slide, Inches(0.3), Inches(5.05), Inches(12.7), Inches(1.1), NAVY)
    add_text_box(slide, "Memory & Streaming",
                 Inches(0.5), Inches(5.1), Inches(12.0), Inches(0.38),
                 font_size=12, bold=True, color=GOLD)
    add_text_box(slide,
                 "Sliding-window deque (maxlen=40 messages = 20 turns) — no external library needed    •    "
                 "stream_chat() resolves all tool calls silently then streams final answer token-by-token via llm.astream()",
                 Inches(0.5), Inches(5.5), Inches(12.0), Inches(0.55),
                 font_size=11, color=LIGHT)


# ── Slide 10: Interfaces ──────────────────────────────────────────────────────

def slide_interfaces(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "User Interfaces", "Three complete interfaces + REST API")

    ifaces = [
        (NAVY, "Rich CLI",
         ["Run: python -m forexmind.main cli",
          "Streaming response display",
          "Commands: /signal, /account, /pairs,\n  /sessions, /backtest, /clear",
          "Colour-coded signal panels",
          "Session status banner"]),
        (BLUE, "FastAPI Web Dashboard",
         ["Run: python -m forexmind.main web",
          "http://localhost:8000",
          "Dark-theme live signal grid",
          "WebSocket: 60s auto-refresh",
          "Embedded Claude chat panel",
          "REST API: /api/signal/{pair}"]),
        (TEAL, "Telegram Bot",
         ["Run: python -m forexmind.main telegram",
          "Commands: /signal PAIR, /signals,\n  /account, /sessions",
          "Inline keyboard buttons",
          "Free text → Claude agent",
          "HTML-formatted messages"]),
    ]

    for i, (col, title, bullets) in enumerate(ifaces):
        lx = Inches(0.3 + i * 4.35)
        add_rect(slide, lx, Inches(1.2), Inches(4.1), Inches(4.9), col)
        add_text_box(slide, title, lx + Inches(0.1), Inches(1.28),
                     Inches(3.9), Inches(0.5),
                     font_size=16, bold=True, color=WHITE)
        for j, b in enumerate(bullets):
            add_text_box(slide, f"• {b}", lx + Inches(0.15), Inches(1.85 + j * 0.56),
                         Inches(3.8), Inches(0.52),
                         font_size=10.5, color=WHITE)


# ── Slide 11: Backtest & Config ───────────────────────────────────────────────

def slide_backtest_config(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Backtesting, Configuration & Security", "Validate before you trade")

    add_bullet_box(slide, "Three-Layer Validation",
        ["Event-driven backtest: bar-by-bar with slippage + commission",
         "Walk-forward CV: TimeSeriesSplit(5), no look-ahead bias",
         "Monte Carlo (1000 runs): shuffle trade order → VaR",
         "Metrics: Win Rate, Profit Factor, Sharpe, Sortino, Max DD"],
        Inches(0.3), Inches(1.2), Inches(6.1), Inches(2.8))

    add_bullet_box(slide, "Configuration System",
        [".env  →  API keys + PAPER_TRADING flag",
         "config.yaml  →  all runtime tunables",
         "Settings dataclass: @lru_cache singleton",
         "Pairs, indicator periods, risk limits, ensemble weights",
         "Retrain interval, news window, session times — all YAML"],
        Inches(6.7), Inches(1.2), Inches(6.3), Inches(2.8))

    add_bullet_box(slide, "Security Model",
        ["No secrets in code — .env excluded from git",
         "Paper trading ON by default; live requires explicit .env opt-in",
         "Kill-switch: hard −5%/day limit halts everything",
         "Parameterized SQL queries (no injection risk)",
         "Telegram responds only to configured TELEGRAM_CHAT_ID"],
        Inches(0.3), Inches(4.2), Inches(6.1), Inches(2.8), bg_color=BLUE)

    add_bullet_box(slide, "Required API Keys",
        ["OANDA — free practice account (api-fxpractice.oanda.com)",
         "ANTHROPIC_API_KEY — console.anthropic.com",
         "ALPHA_VANTAGE_API_KEY — alphavantage.co (free 25/day)",
         "FINNHUB_API_KEY — finnhub.io (free tier)",
         "TELEGRAM_BOT_TOKEN + CHAT_ID — optional"],
        Inches(6.7), Inches(4.2), Inches(6.3), Inches(2.8), bg_color=BLUE)


# ── Slide 12: Performance Targets ────────────────────────────────────────────

def slide_performance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, WHITE)
    title_bar(slide, "Performance Targets & Getting Started", "What success looks like")

    metrics = [
        ("≥70%", "Directional\nAccuracy"),
        ("≥1.5",  "Profit\nFactor"),
        ("≥1.0",  "Sharpe\nRatio"),
        ("≥2:1",  "R:R per\nTrade"),
        ("≤8%",   "Max Monthly\nDrawdown"),
        ("<2s",   "Signal\nLatency"),
    ]

    for i, (val, label) in enumerate(metrics):
        row = i // 3
        col_pos = i % 3
        lx = Inches(0.3 + col_pos * 4.35)
        ty = Inches(1.2 + row * 1.65)
        color = GREEN if i < 4 else GOLD
        add_rect(slide, lx, ty, Inches(4.1), Inches(1.45), color)
        add_text_box(slide, val, lx, ty + Inches(0.1),
                     Inches(4.1), Inches(0.75),
                     font_size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, label, lx, ty + Inches(0.85),
                     Inches(4.1), Inches(0.5),
                     font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

    # Getting started steps
    add_rect(slide, Inches(0.3), Inches(4.65), Inches(12.7), Inches(1.8), NAVY)
    add_text_box(slide, "Getting Started (3 steps)",
                 Inches(0.5), Inches(4.7), Inches(12.0), Inches(0.38),
                 font_size=12, bold=True, color=GOLD)
    steps = ("1. Add API keys to /home/wilson/Forex/.env   "
             "→   2. Train ML models: python -m forexmind.main train EUR_USD   "
             "→   3. Launch: python -m forexmind.main cli  |  web  |  telegram")
    add_text_box(slide, steps, Inches(0.5), Inches(5.1), Inches(12.0), Inches(1.0),
                 font_size=11, color=LIGHT)


# ── Slide 13: Thank You / Q&A ─────────────────────────────────────────────────

def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, NAVY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), GOLD)
    add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), GOLD)

    add_text_box(slide, "ForexMind", Inches(1), Inches(1.5), Inches(11.33), Inches(1.4),
                 font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, "Ready to Trade Smarter",
                 Inches(1), Inches(2.8), Inches(11.33), Inches(0.7),
                 font_size=24, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_text_box(slide, "Questions & Discussion",
                 Inches(1), Inches(3.5), Inches(11.33), Inches(0.55),
                 font_size=18, color=LIGHT, align=PP_ALIGN.CENTER)

    links = [
        "Project: /home/wilson/Forex/forexmind/",
        "Tests:  python -m pytest forexmind/tests/ -v   (35/35 passing)",
        "Docs:   /home/wilson/Forex/docs/",
    ]
    for i, link in enumerate(links):
        add_text_box(slide, link, Inches(2), Inches(4.3 + i * 0.55), Inches(9.33), Inches(0.45),
                     font_size=13, color=MID, align=PP_ALIGN.CENTER)


# ── Main build ────────────────────────────────────────────────────────────────

def build_pptx():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_agenda(prs)
    slide_overview(prs)
    slide_architecture(prs)
    slide_data_layer(prs)
    slide_indicators(prs)
    slide_ensemble(prs)
    slide_risk(prs)
    slide_agent(prs)
    slide_interfaces(prs)
    slide_backtest_config(prs)
    slide_performance(prs)
    slide_thank_you(prs)

    prs.save(OUTPUT)
    print(f"✓ PowerPoint written → {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    build_pptx()
